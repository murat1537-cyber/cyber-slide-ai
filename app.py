import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE # Yeni ekledik: Şekiller için
import json
import io

# --- YARDIMCI FONKSİYONLAR ---
def hex_to_rgb(hex_str):
    """HEX renk kodunu RGB'ye çevirir."""
    hex_str = hex_str.lstrip('#')
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def create_pptx(json_data):
    """Gelen JSON verisini sıfırdan PPTX dosyasına çizer."""
    data = json.loads(json_data)
    prs = Presentation()
    
    bg_rgb = hex_to_rgb(data["presentation_metadata"]["global_background_color_hex"])
    accent_rgb = hex_to_rgb(data["presentation_metadata"]["global_accent_color_hex"])
    
    for slide_data in data["slides"]:
        # Boş şablon
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Arka Planı Boya
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(bg_rgb[0], bg_rgb[1], bg_rgb[2])
        
        # Başlık Ekle
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data["slide_title"]
        title_p.font.bold = True
        title_p.font.size = Pt(36)
        title_p.font.color.rgb = RGBColor(accent_rgb[0], accent_rgb[1], accent_rgb[2])
        
        # --- İÇERİK VE GÖRSEL MİZANPAJ (LAYOUT) ---
        layout_type = slide_data.get("layout_type", "text_only")
        
        if layout_type == "text_only":
            # Normal metin kutusu
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8.5), Inches(5))
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            
            for bullet in slide_data["content_bullets"]:
                p = body_frame.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(22)
                text_color = 255 if sum(bg_rgb) < 380 else 0
                p.font.color.rgb = RGBColor(text_color, text_color, text_color)
                
        elif layout_type == "text_and_diagram":
            # Metin solda, şema sağda
            
            # Sol Taraf: Metin Kutusu
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            for bullet in slide_data["content_bullets"]:
                p = text_frame.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                text_color = 255 if sum(bg_rgb) < 380 else 0
                p.font.color.rgb = RGBColor(text_color, text_color, text_color)
                
            # Sağ Taraf: Şema/Şekil
            diagram_instruction = slide_data.get("visual_element", {}).get("type", "")
            
            if diagram_instruction == "simple_flowchart":
                # Yapay zekanın "akış şeması" isteğini temel şekillerle (kutular ve oklar) çiziyoruz.
                shapes_to_draw = slide_data.get("visual_element", {}).get("shapes", [])
                
                # Örnek: Basit bir dikdörtgen ekleyip boyayalım (Daha karmaşık çizimler için fonksiyonu geliştirebiliriz)
                for shape in shapes_to_draw:
                    shape_type_str = shape.get("type", "")
                    if shape_type_str == "rectangle":
                        new_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, 
                            Inches(shape.get("x_inches", 5)), 
                            Inches(shape.get("y_inches", 2)), 
                            Inches(shape.get("width_inches", 2)), 
                            Inches(shape.get("height_inches", 1))
                        )
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = RGBColor(accent_rgb[0], accent_rgb[1], accent_rgb[2])
                        new_shape.text_frame.text = shape.get("text", "")
                        new_shape.text_frame.paragraphs[0].font.size = Pt(16)
                        
            elif diagram_instruction == "network_topology":
                # Yapay zeka "Ağ Topolojisi" dediğinde, arka planda bir görsel üreten API'yi 
                # (örneğin DALL-E) arayıp resmi buraya koyabiliriz.
                st.warning("Ağ topolojisi şeması oluşturmak için DALL-E API entegrasyonu gereklidir. Şimdilik bu kısmı atlıyoruz.")
                pass

    # Dosyayı hafızada tut, diske kaydetme
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- UYGULAMA ARAYÜZÜ (STREAMLIT) ---
st.set_page_config(page_title="Cyber-Slide AI", page_icon="🛡️", layout="wide")

st.title("🛡️ Cyber-Slide AI: Akıllı Sunum Mimarı")
st.markdown("Siber güvenlik eğitimleriniz için saniyeler içinde profesyonel ve özelleştirilmiş slaytlar oluşturun.")

# API Anahtarını Streamlit Secrets'tan al
API_KEY = st.secrets.get("GEMINI_API_KEY", "")

if not API_KEY:
    st.error("Lütfen Streamlit Secrets bölümüne GEMINI_API_KEY ekleyin!")
    st.stop()

genai.configure(api_key=API_KEY)

# Sol Menü (Ayarlar)
with st.sidebar:
    st.header("⚙️ Eğitim Parametreleri")
    topic = st.text_input("Konu Nedir?", "Zero Trust Architecture in Cloud")
    language = st.selectbox("Sunum Dili", ["English", "Nederlands (Dutch)"])
    slide_count = st.slider("Slayt Sayısı", min_value=3, max_value=20, value=5)
    design_prompt = st.text_area("Tasarım ve Renk Tercihi", "Koyu gri arka plan, başlıklar siber güvenlik hissi veren neon yeşil olsun. Ciddi ve kurumsal.")

# Ana Ekran - Üretim Butonu
if st.button("🚀 Sunumu Üret", type="primary"):
    with st.spinner("Yapay zeka konuyu analiz ediyor, görselleri tasarlıyor ve slaytları sıfırdan çiziyor... Lütfen bekleyin."):
        
        # Gemini Modeli Ayarları (JSON formatında çıkış vermeye zorluyoruz)
        # GÜNCELLEME: En kararlı sürüm olan gemini-1.5-flash'ı kullanıyoruz.
        model = genai.GenerativeModel('gemini-2.5-flash', generation_config={"response_mime_type": "application/json"})
        
        system_prompt = f"""
        You are an elite Cybersecurity Instructor. Generate a professional presentation.
        Topic: {topic}
        Language: {language} (Keep strict cybersecurity terms in English).
        Slide Count: Exactly {slide_count} slides.
        Design Vibe: {design_prompt}
        
        You must decide the exact color HEX codes based on the Design Vibe.
        For each slide, you must choose a "layout_type" ("text_only" or "text_and_diagram").
        If "text_and_diagram" is chosen, provide a "simple_flowchart" or "network_topology" specification under "visual_element".
        
        Output EXACTLY in this JSON structure:
        {{
          "presentation_metadata": {{
            "global_background_color_hex": "#HEXCODE",
            "global_accent_color_hex": "#HEXCODE"
          }},
          "slides": [
            {{
              "slide_number": 1,
              "layout_type": "text_and_diagram",
              "slide_title": "Slide Title Here",
              "content_bullets": ["Bullet 1", "Bullet 2"],
              "visual_element": {{
                "type": "simple_flowchart",
                "shapes": [
                  {{
                    "type": "rectangle",
                    "x_inches": 5.0,
                    "y_inches": 2.0,
                    "width_inches": 2.0,
                    "height_inches": 1.0,
                    "text": "User"
                  }}
                ]
              }}
            }}
          ]
        }}
        """
        
        try:
            # Gemini'den veriyi al
            response = model.generate_content(system_prompt)
            json_result = response.text
            
            # PPTX Dosyasını Çiz
            ppt_file = create_pptx(json_result)
            
            st.success("🎉 Sunumunuz başarıyla hazırlandı!")
            
            # İndirme Butonu
            st.download_button(
                label="📥 PowerPoint (.pptx) Dosyasını İndir",
                data=ppt_file,
                file_name=f"{topic.replace(' ', '_')}_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
        except Exception as e:
            st.error(f"Bir hata oluştu: {e}")
